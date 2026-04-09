import os
import uuid
import logging
from flask import Blueprint, request, jsonify, current_app, send_from_directory
from werkzeug.utils import secure_filename
from app import db
from app.models.slide import Slide, SlidePage
from app.models.course import Course
from app.models.knowledge_point import KnowledgePoint
from app.models.quiz import Quiz, QuizAttempt
from app.auth_utils import require_teacher

slides_bp = Blueprint("slides", __name__)
logger = logging.getLogger(__name__)

ALLOWED_EXTENSIONS = {"pdf", "ppt", "pptx"}


def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def get_file_extension(filename):
    _, ext = os.path.splitext(filename)
    return ext.lstrip(".").lower()


@slides_bp.route("/upload", methods=["POST"])
def upload_slide():
    forbidden = require_teacher()
    if forbidden:
        return forbidden

    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected"}), 400

    if not allowed_file(file.filename):
        return jsonify({"error": "Unsupported file type. Use PDF, PPT, or PPTX."}), 400

    course_id = request.form.get("course_id")
    if not course_id:
        return jsonify({"error": "course_id is required"}), 400

    course = db.session.get(Course, int(course_id))
    if not course:
        return jsonify({"error": "Course not found"}), 404

    original_filename = secure_filename(file.filename)
    ext = get_file_extension(file.filename)
    if not ext:
        return jsonify({"error": "Invalid file name. Please use a .pdf, .ppt, or .pptx file."}), 400
    unique_filename = f"{uuid.uuid4().hex}.{ext}"
    upload_dir = os.path.join(current_app.config["UPLOAD_FOLDER"], "slides")
    file_path = os.path.join(upload_dir, unique_filename)
    file.save(file_path)

    slide = Slide(
        course_id=int(course_id),
        filename=unique_filename,
        original_filename=file.filename,
        file_type=ext,
        file_path=file_path,
    )
    db.session.add(slide)
    db.session.commit()

    # Extract text from PDF (basic Phase 1 processing)
    if ext == "pdf":
        try:
            _process_pdf(slide)
        except Exception as e:
            logger.warning("PDF processing error for slide %s: %s", slide.id, e)
    elif ext in ("ppt", "pptx"):
        try:
            _process_pptx(slide)
        except Exception as e:
            logger.warning("PPTX processing error for slide %s: %s", slide.id, e)

    # Generate embeddings for slide pages (async-friendly)
    try:
        from app.services.alignment_service import embed_slide_pages
        embed_slide_pages(slide.id)
    except Exception as e:
        logger.warning("Slide embedding error: %s", e)

    return jsonify(slide.to_dict()), 201


def _process_pdf(slide):
    import fitz  # PyMuPDF

    doc = fitz.open(slide.file_path)
    slide.total_pages = len(doc)

    thumbs_dir = os.path.join(
        os.path.dirname(slide.file_path), "thumbnails", str(slide.id)
    )
    os.makedirs(thumbs_dir, exist_ok=True)

    for i, page in enumerate(doc):
        text = page.get_text() or ""
        # Render page to image (2x zoom for clarity)
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat)
        img_filename = f"page_{i + 1}.png"
        img_path = os.path.join(thumbs_dir, img_filename)
        pix.save(img_path)

        slide_page = SlidePage(
            slide_id=slide.id,
            page_number=i + 1,
            content_text=text,
            thumbnail_path=f"thumbnails/{slide.id}/{img_filename}",
        )
        db.session.add(slide_page)

    doc.close()
    slide.processed = True
    db.session.commit()


def _process_pptx(slide):
    """Extract text and render images from PPT/PPTX files."""
    import subprocess
    import fitz  # PyMuPDF
    from pptx import Presentation

    prs = Presentation(slide.file_path)
    slide.total_pages = len(prs.slides)

    # Collect text per slide first
    slide_texts = []
    for pptx_slide in prs.slides:
        texts = []
        for shape in pptx_slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        texts.append(para_text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append(" | ".join(row_texts))
        slide_texts.append("\n".join(texts))

    thumbs_dir = os.path.join(
        os.path.dirname(slide.file_path), "thumbnails", str(slide.id)
    )
    os.makedirs(thumbs_dir, exist_ok=True)

    # Convert PPTX to PDF using LibreOffice for image rendering
    pdf_path = None
    try:
        result = subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                thumbs_dir,
                slide.file_path,
            ],
            capture_output=True,
            timeout=120,
        )
        base_name = os.path.splitext(os.path.basename(slide.file_path))[0]
        candidate = os.path.join(thumbs_dir, f"{base_name}.pdf")
        if os.path.exists(candidate):
            pdf_path = candidate
    except Exception as e:
        logger.warning("LibreOffice conversion failed: %s", e)

    if pdf_path:
        doc = fitz.open(pdf_path)
        for i in range(min(len(doc), slide.total_pages)):
            page = doc[i]
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat)
            img_filename = f"page_{i + 1}.png"
            img_path = os.path.join(thumbs_dir, img_filename)
            pix.save(img_path)

            slide_page = SlidePage(
                slide_id=slide.id,
                page_number=i + 1,
                content_text=slide_texts[i] if i < len(slide_texts) else "",
                thumbnail_path=f"thumbnails/{slide.id}/{img_filename}",
            )
            db.session.add(slide_page)
        doc.close()
        os.remove(pdf_path)
    else:
        # Fallback: text-only if LibreOffice unavailable
        for i, text in enumerate(slide_texts):
            slide_page = SlidePage(
                slide_id=slide.id,
                page_number=i + 1,
                content_text=text,
            )
            db.session.add(slide_page)

    slide.processed = True
    db.session.commit()


@slides_bp.route("/course/<int:course_id>", methods=["GET"])
def get_slides_by_course(course_id):
    slides = Slide.query.filter_by(course_id=course_id).order_by(Slide.created_at.desc()).all()
    return jsonify([s.to_dict() for s in slides])


@slides_bp.route("/<int:slide_id>", methods=["GET"])
def get_slide(slide_id):
    slide = db.session.get(Slide, slide_id)
    if not slide:
        return jsonify({"error": "Slide not found"}), 404
    return jsonify(slide.to_dict())


@slides_bp.route("/file/<path:filename>", methods=["GET"])
def serve_slide_file(filename):
    upload_dir = os.path.join(current_app.config["UPLOAD_FOLDER"], "slides")
    return send_from_directory(upload_dir, filename)


@slides_bp.route("/page-image/<int:page_id>", methods=["GET"])
def serve_page_image(page_id):
    """Serve the rendered image for a slide page."""
    page = db.session.get(SlidePage, page_id)
    if not page or not page.thumbnail_path:
        return jsonify({"error": "Page image not found"}), 404
    upload_dir = os.path.join(current_app.config["UPLOAD_FOLDER"], "slides")
    return send_from_directory(upload_dir, page.thumbnail_path)


@slides_bp.route("/<int:slide_id>", methods=["DELETE"])
def delete_slide(slide_id):
    forbidden = require_teacher()
    if forbidden:
        return forbidden

    slide = db.session.get(Slide, slide_id)
    if not slide:
        return jsonify({"error": "Slide not found"}), 404

    # Get all slide pages associated with this slide
    slide_pages = SlidePage.query.filter_by(slide_id=slide_id).all()
    
    # 1. Delete QuizAttempts that reference Quizzes related to this slide's knowledge points
    for page in slide_pages:
        kps = KnowledgePoint.query.filter_by(slide_page_id=page.id).all()
        for kp in kps:
            quizzes = Quiz.query.filter_by(knowledge_point_id=kp.id).all()
            for quiz in quizzes:
                QuizAttempt.query.filter_by(quiz_id=quiz.id).delete()
    
    # 2. Delete Quizzes
    for page in slide_pages:
        kps = KnowledgePoint.query.filter_by(slide_page_id=page.id).all()
        for kp in kps:
            Quiz.query.filter_by(knowledge_point_id=kp.id).delete()
    
    # 3. Delete KnowledgePoints
    for page in slide_pages:
        KnowledgePoint.query.filter_by(slide_page_id=page.id).delete()
    
    # 4. Clean up files
    if os.path.exists(slide.file_path):
        os.remove(slide.file_path)

    # Clean up thumbnail directory
    import shutil
    thumbs_dir = os.path.join(
        os.path.dirname(slide.file_path), "thumbnails", str(slide.id)
    )
    if os.path.isdir(thumbs_dir):
        shutil.rmtree(thumbs_dir, ignore_errors=True)

    # 5. Delete Slide (SlidePage will be cascade deleted)
    db.session.delete(slide)
    db.session.commit()
    return jsonify({"message": "Slide deleted"})
